from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
import io
import tempfile
import os

def create_presentation(analysis_results, insights, recommendations):
    """
    Create a PowerPoint presentation with the analysis results
    """
    # Create a new presentation
    prs = Presentation()
    
    # Create title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "E-commerce Funnel Analysis"
    subtitle.text = "Identifying Conversion Issues and Improvement Strategies"
    
    # Add an overview slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Overview"
    tf = body.text_frame
    tf.text = "Analysis of the e-commerce conversion funnel:"
    
    p = tf.add_paragraph()
    p.text = "• Home → Search → Payment → Confirmation"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"• Overall conversion rate: {analysis_results['overall']['conversion_rate']}%"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"• {analysis_results['user_counts']['total']} total users analyzed"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"• New users: {analysis_results['user_counts']['new']} ({round(analysis_results['user_counts']['new']/analysis_results['user_counts']['total']*100, 1)}%)"
    p.level = 1
    
    # Add a funnel analysis slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Funnel Analysis"
    tf = body.text_frame
    
    funnel_df = analysis_results['overall']['funnel']
    # Create a temporary file for the funnel chart
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        funnel_data = [
            funnel_df.loc[0, 'Users'],
            funnel_df.loc[1, 'Users'],
            funnel_df.loc[2, 'Users'],
            funnel_df.loc[3, 'Users']
        ]
        stages = funnel_df['Stage'].tolist()
        
        # Create the funnel chart using matplotlib
        plt.figure(figsize=(8, 6))
        plt.bar(stages, funnel_data, color=['#0068c9', '#83c9ff', '#29b09d', '#7defa1'])
        for i, v in enumerate(funnel_data):
            plt.text(i, v + 0.1, str(v), ha='center')
        plt.title('User Funnel')
        plt.tight_layout()
        plt.savefig(tmp.name)
        plt.close()
        
        # Add the chart to the slide
        slide.shapes.add_picture(tmp.name, Inches(1), Inches(2), width=Inches(8))
    
    # Clean up the temporary file
    os.unlink(tmp.name)
    
    # Add a conversion rates slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Conversion Rates Between Stages"
    tf = body.text_frame
    
    # Display conversion rates
    p = tf.add_paragraph()
    p.text = f"• Home to Search: {funnel_df.loc[1, 'Conversion_Rate']}%"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = f"• Search to Payment: {funnel_df.loc[2, 'Conversion_Rate']}%"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = f"• Payment to Confirmation: {funnel_df.loc[3, 'Conversion_Rate']}%"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = f"• Overall (Home to Confirmation): {analysis_results['overall']['conversion_rate']}%"
    p.level = 0
    
    # Add a drop-off analysis slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Drop-off Analysis"
    tf = body.text_frame
    
    drop_off_df = analysis_results['overall']['drop_off']
    
    # Display drop-off rates
    for i in range(len(drop_off_df)):
        p = tf.add_paragraph()
        p.text = f"• {drop_off_df.loc[i, 'Stage']}: {drop_off_df.loc[i, 'Drop_Off_Percentage']}% ({drop_off_df.loc[i, 'Drop_Off_Count']} users)"
        p.level = 0
    
    # Add a segment comparison slide - Device
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Device Comparison"
    tf = body.text_frame
    
    device_segments = analysis_results['segments']['device']
    for device, data in device_segments.items():
        p = tf.add_paragraph()
        p.text = f"• {device}: {data['overall_conversion']}% overall conversion"
        p.level = 0
        
        funnel = data['funnel_df']
        p = tf.add_paragraph()
        p.text = f"  - Home to Search: {funnel.loc[1, 'Conversion_Rate']}%"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"  - Search to Payment: {funnel.loc[2, 'Conversion_Rate']}%"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"  - Payment to Confirmation: {funnel.loc[3, 'Conversion_Rate']}%"
        p.level = 1
    
    # Add a segment comparison slide - Gender
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Gender Comparison"
    tf = body.text_frame
    
    gender_segments = analysis_results['segments']['gender']
    for gender, data in gender_segments.items():
        p = tf.add_paragraph()
        p.text = f"• {gender}: {data['overall_conversion']}% overall conversion"
        p.level = 0
        
        funnel = data['funnel_df']
        p = tf.add_paragraph()
        p.text = f"  - Home to Search: {funnel.loc[1, 'Conversion_Rate']}%"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"  - Search to Payment: {funnel.loc[2, 'Conversion_Rate']}%"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"  - Payment to Confirmation: {funnel.loc[3, 'Conversion_Rate']}%"
        p.level = 1
    
    # Add new vs existing users comparison slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "New vs Existing Users"
    tf = body.text_frame
    
    new_data = analysis_results['segments']['user_type']['new']
    existing_data = analysis_results['segments']['user_type']['existing']
    
    p = tf.add_paragraph()
    p.text = f"• New Users: {new_data['overall_conversion']}% overall conversion"
    p.level = 0
    
    new_funnel = new_data['funnel']
    p = tf.add_paragraph()
    p.text = f"  - Home to Search: {new_funnel.loc[1, 'Conversion_Rate']}%"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"  - Search to Payment: {new_funnel.loc[2, 'Conversion_Rate']}%"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"  - Payment to Confirmation: {new_funnel.loc[3, 'Conversion_Rate']}%"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"• Existing Users: {existing_data['overall_conversion']}% overall conversion"
    p.level = 0
    
    existing_funnel = existing_data['funnel']
    p = tf.add_paragraph()
    p.text = f"  - Home to Search: {existing_funnel.loc[1, 'Conversion_Rate']}%"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"  - Search to Payment: {existing_funnel.loc[2, 'Conversion_Rate']}%"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"  - Payment to Confirmation: {existing_funnel.loc[3, 'Conversion_Rate']}%"
    p.level = 1
    
    # Add key insights slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Key Insights"
    tf = body.text_frame
    
    for insight in insights:
        p = tf.add_paragraph()
        p.text = f"• {insight}"
        p.level = 0
    
    # Add recommendations slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Strategic Recommendations"
    tf = body.text_frame
    
    for recommendation in recommendations[:7]:  # Limit to first 7 recommendations
        p = tf.add_paragraph()
        p.text = f"• {recommendation}"
        p.level = 0
    
    # Add additional recommendations slide if needed
    if len(recommendations) > 7:
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        
        title.text = "Strategic Recommendations (Continued)"
        tf = body.text_frame
        
        for recommendation in recommendations[7:]:
            p = tf.add_paragraph()
            p.text = f"• {recommendation}"
            p.level = 0
    
    # Add a conclusion slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Conclusion"
    tf = body.text_frame
    tf.text = "Key Actions to Improve Conversion:"
    
    p = tf.add_paragraph()
    p.text = "• Focus on optimizing the biggest drop-off point in the funnel"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Implement specific strategies for new users to improve their conversion"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Consider device-specific optimizations, especially for mobile users"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Implement A/B testing to continuously improve the conversion funnel"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Set up a monitoring system to track improvements over time"
    p.level = 1
    
    # Save the presentation to a BytesIO object
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    
    return ppt_bytes
